#!/usr/bin/env python3
"""
generate-ppt.py — report-slides.md의 slide_meta를 파싱하여 PPTX 생성

사용법:
    python3 scripts/generate-ppt.py <project-name>

예시:
    python3 scripts/generate-ppt.py my-research

필수 패키지:
    pip3 install python-pptx

입력: {project}/reports/report-slides.md (slide_meta HTML 코멘트 포함)
출력: {project}/reports/report-slides.pptx
"""

import sys
import re
import os
import yaml
from pathlib import Path

def checkDependency():
    """python-pptx 설치 확인"""
    try:
        from pptx import Presentation
        return True
    except ImportError:
        print("❌ python-pptx가 설치되어 있지 않습니다.")
        print("   설치: pip3 install python-pptx")
        return False


# ============================================================
# 차트 연동 함수
# ============================================================

def loadChartIndex(projectPath):
    """chart-index.yaml 로드. 없으면 빈 리스트 반환."""
    chartIndexPath = Path(projectPath) / "reports" / "chart-index.yaml"
    if not chartIndexPath.exists():
        return []
    try:
        with open(chartIndexPath, 'r', encoding='utf-8') as f:
            data = yaml.safe_load(f)
        if data and 'charts' in data:
            return data['charts']
    except Exception as e:
        print(f"  [경고] chart-index.yaml 로드 실패: {e}")
    return []


def findChartsForSlide(slideClaims, charts):
    """슬라이드의 source_claims와 매칭되는 차트 찾기.
    slideClaims: 슬라이드에 연결된 Claim ID 리스트
    charts: chart-index.yaml에서 로드한 차트 목록
    반환: 매칭된 차트 엔트리 리스트
    """
    if not slideClaims or not charts:
        return []
    matched = []
    for chart in charts:
        relatedClaims = chart.get('related_claims', [])
        if any(claim in relatedClaims for claim in slideClaims):
            matched.append(chart)
    return matched


def validateActionTitle(title):
    """Action Title 기본 검증 — 주제형 타이틀 경고"""
    # 주제형 패턴: 명사형 종결
    themePatterns = ['분석', '현황', '전망', '비교', '요약', '개요', '검토', '조사']
    for p in themePatterns:
        if title.strip().endswith(p):
            return False, f"주제형 타이틀 감지: '{title}' ('{p}'로 종결)"
    if len(title.strip()) < 5:
        return False, f"타이틀 너무 짧음: '{title}'"
    return True, None


def addChartToSlide(slide, chartPath, chartType="standard", position="right"):
    """슬라이드에 차트 이미지를 추가.
    chartType: 차트 유형에 따라 레이아웃 자동 조정
    position: "right" = 우측 배치, "bottom" = 하단 배치
    """
    from pptx.util import Inches

    imgPath = Path(chartPath)
    if not imgPath.exists():
        return False

    # 차트 유형별 레이아웃 설정
    layoutConfig = {
        "marimekko": {"x": Inches(1.5), "y": Inches(3.0), "w": Inches(10.0), "h": Inches(3.5)},
        "harveyball": {"x": Inches(7.5), "y": Inches(1.5), "w": Inches(5.0), "h": Inches(5.0)},
    }

    try:
        if chartType in layoutConfig:
            # 특수 차트 유형: 전용 레이아웃
            cfg = layoutConfig[chartType]
            slide.shapes.add_picture(str(imgPath), cfg["x"], cfg["y"], cfg["w"], cfg["h"])
        elif position == "right":
            # 우측 배치: 슬라이드 오른쪽 영역에 차트 삽입
            slide.shapes.add_picture(
                str(imgPath), Inches(7.5), Inches(1.5), Inches(5.3), Inches(4.0)
            )
        elif position == "bottom":
            # 하단 배치: 슬라이드 아래쪽 전체 너비에 차트 삽입
            slide.shapes.add_picture(
                str(imgPath), Inches(1.5), Inches(4.2), Inches(10.0), Inches(3.0)
            )
        return True
    except Exception as e:
        print(f"  [경고] 차트 이미지 삽입 실패 ({chartPath}): {e}")
        return False

def parseSlidesMd(filePath):
    """report-slides.md에서 슬라이드 섹션을 파싱"""
    with open(filePath, 'r', encoding='utf-8') as f:
        content = f.read()

    slides = []

    # slide_meta HTML 코멘트가 있으면 파싱
    metaPattern = r'<!-- slide_meta\n(.*?)\n-->'
    metaBlocks = re.findall(metaPattern, content, re.DOTALL)

    if metaBlocks:
        # slide_meta 기반 파싱
        for block in metaBlocks:
            slide = parseMetaBlock(block)
            slides.append(slide)
    else:
        # slide_meta 없으면 ## Slide/슬라이드 N: 패턴으로 폴백 (영어+한국어 대응)
        slidePattern = r'## (?:Slide|슬라이드)\s+(\d+):\s*(.*?)\n(.*?)(?=## (?:Slide|슬라이드|별첨)\s+[\dA-Z]+[:\.]|---|\Z)'
        matches = re.findall(slidePattern, content, re.DOTALL)
        for num, title, body in matches:
            # Action Title 검증
            isValid, msg = validateActionTitle(title.strip())
            if not isValid:
                print(f"  ⚠️ Slide {num}: {msg}")
            slides.append({
                'slide_id': int(num),
                'title': title.strip(),
                'layout': 'content',
                'body': body.strip()
            })

        # 별첨(Appendix) 슬라이드도 파싱
        appendixPattern = r'## 별첨\s+([A-Z])[\.:]\s*(.*?)\n(.*?)(?=## 별첨\s+[A-Z][\.:]+|---|\Z)'
        appendixMatches = re.findall(appendixPattern, content, re.DOTALL)
        for letter, title, body in appendixMatches:
            slides.append({
                'slide_id': 100 + ord(letter) - ord('A'),  # A=100, B=101...
                'title': f'[별첨 {letter}] {title.strip()}',
                'layout': 'content',
                'body': body.strip()
            })

    return slides

def parseMetaBlock(block):
    """slide_meta YAML 블록을 딕셔너리로 변환 (간이 파싱)"""
    slide = {
        'slide_id': 0,
        'title': '',
        'layout': 'content',
        'body': '',
        'source_claims': []
    }

    for line in block.strip().split('\n'):
        line = line.strip()
        if line.startswith('slide_id:'):
            slide['slide_id'] = int(line.split(':', 1)[1].strip())
        elif line.startswith('title:'):
            slide['title'] = line.split(':', 1)[1].strip().strip('"')
        elif line.startswith('layout:'):
            slide['layout'] = line.split(':', 1)[1].strip()

    # source_claims 추출: content_blocks 내 source_claims 리스트를 수집
    claimPattern = r'source_claims:\s*\[([^\]]*)\]'
    matches = re.findall(claimPattern, block)
    allClaims = []
    for match in matches:
        # "MGN-02, FRV-01" → ['MGN-02', 'FRV-01']
        claims = [c.strip().strip('"').strip("'") for c in match.split(',') if c.strip()]
        allClaims.extend(claims)
    # 중복 제거 (순서 유지)
    seen = set()
    uniqueClaims = []
    for c in allClaims:
        if c not in seen:
            seen.add(c)
            uniqueClaims.append(c)
    slide['source_claims'] = uniqueClaims

    return slide

def mdToPlaintext(mdText):
    """마크다운을 간이 플레인텍스트로 변환"""
    # 테이블 헤더 구분선 제거
    text = re.sub(r'\|[-:]+\|[-:|\s]+\|', '', mdText)
    # 볼드/이탤릭 제거
    text = re.sub(r'\*\*(.*?)\*\*', r'\1', text)
    text = re.sub(r'\*(.*?)\*', r'\1', text)
    # 헤더 마커 제거
    text = re.sub(r'^#{1,4}\s+', '', text, flags=re.MULTILINE)
    # 빈 줄 정리
    text = re.sub(r'\n{3,}', '\n\n', text)
    return text.strip()

def generatePptx(slides, outputPath, charts=None, projectPath=None):
    """파싱된 슬라이드 데이터로 PPTX 생성.
    charts: chart-index.yaml에서 로드한 차트 목록 (없으면 None)
    projectPath: 프로젝트 루트 경로 (차트 파일 절대경로 계산용)
    """
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN

    prs = Presentation()
    prs.slide_width = Inches(13.333)  # 와이드스크린 16:9
    prs.slide_height = Inches(7.5)

    chartInsertCount = 0

    for slideData in slides:
        layout = slideData.get('layout', 'content')
        title = slideData.get('title', '')
        body = slideData.get('body', '')
        sourceClaims = slideData.get('source_claims', [])

        # 이 슬라이드에 매칭되는 차트 탐색
        matchedCharts = findChartsForSlide(sourceClaims, charts) if charts else []
        hasChart = len(matchedCharts) > 0

        if layout == 'title_slide':
            # 타이틀 슬라이드
            slideLayout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(slideLayout)
            slide.shapes.title.text = title
            if slide.placeholders[1]:
                slide.placeholders[1].text = mdToPlaintext(body)[:200]
        else:
            # 일반 콘텐츠 슬라이드
            slideLayout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(slideLayout)
            slide.shapes.title.text = title

            # 본문 콘텐츠
            if body and len(slide.placeholders) > 1:
                tf = slide.placeholders[1].text_frame
                tf.clear()
                plainBody = mdToPlaintext(body)

                # 차트가 있으면 텍스트 영역을 좁혀서 좌측에 배치
                maxLines = 25 if hasChart else 30

                # 줄 단위로 추가
                lines = plainBody.split('\n')
                for i, line in enumerate(lines[:maxLines]):
                    if i == 0:
                        tf.text = line
                    else:
                        p = tf.add_paragraph()
                        p.text = line

                    # 폰트 설정
                    for paragraph in tf.paragraphs:
                        paragraph.font.size = Pt(14)

            # 차트 이미지 삽입
            if hasChart and projectPath:
                # 첫 번째 차트는 우측, 추가 차트가 있으면 하단
                for idx, chart in enumerate(matchedCharts[:2]):  # 슬라이드당 최대 2개
                    chartFilePath = Path(projectPath) / chart.get('file', '')
                    position = "right" if idx == 0 else "bottom"
                    if addChartToSlide(slide, chartFilePath, position):
                        chartInsertCount += 1

    prs.save(outputPath)
    return len(slides), chartInsertCount

def main():
    if len(sys.argv) < 2:
        print("사용법: python3 scripts/generate-ppt.py <project-name>")
        sys.exit(1)

    project = sys.argv[1]
    repoDir = Path(__file__).resolve().parent.parent
    projectPath = repoDir / project
    slidesPath = projectPath / "reports" / "report-slides.md"
    outputPath = projectPath / "reports" / "report-slides.pptx"

    # 의존성 확인
    if not checkDependency():
        sys.exit(1)

    # 입력 파일 확인
    if not slidesPath.exists():
        print(f"❌ {slidesPath} 파일이 없습니다.")
        print(f"   Phase 4 보고서 생성을 먼저 완료하세요.")
        sys.exit(1)

    print(f"=== PPT 생성 시작 ===")
    print(f"입력: {slidesPath}")
    print(f"출력: {outputPath}")

    # chart-index.yaml 로드 (없으면 빈 리스트 — 하위 호환)
    charts = loadChartIndex(str(projectPath))
    if charts:
        print(f"차트: {len(charts)}개 로드됨 (chart-index.yaml)")
    else:
        print(f"차트: chart-index.yaml 없음 — 텍스트 전용 PPT 생성")
    print()

    # 슬라이드 파싱
    slides = parseSlidesMd(slidesPath)
    if not slides:
        print("❌ 슬라이드를 파싱할 수 없습니다.")
        print("   report-slides.md에 '## 슬라이드 N:' 패턴이 있는지 확인하세요.")
        sys.exit(1)

    print(f"파싱된 슬라이드: {len(slides)}개")
    for s in slides:
        claimInfo = f" (claims: {', '.join(s['source_claims'])})" if s.get('source_claims') else ""
        print(f"  [{s['slide_id']}] {s['title'][:50]}{claimInfo}")

    # PPTX 생성 (차트 연동 포함)
    slideCount, chartCount = generatePptx(
        slides, str(outputPath),
        charts=charts if charts else None,
        projectPath=str(projectPath)
    )

    print()
    print(f"=== PPT 생성 완료 ===")
    print(f"  {slideCount}개 슬라이드 → {outputPath}")
    if chartCount > 0:
        print(f"  {chartCount}개 차트 이미지 삽입됨")
    print()
    print(f"열기: open {outputPath}")

if __name__ == "__main__":
    main()
